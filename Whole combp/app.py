from flask import Flask, render_template, request, redirect, url_for, flash, jsonify
from werkzeug.utils import secure_filename
import os
import anthropic
from pptx import Presentation
import PyPDF2
import json
from datetime import datetime

app = Flask(__name__)
app.secret_key = 'your-secret-key-change-this'
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['ALLOWED_EXTENSIONS'] = {'pdf', 'pptx'}

os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

#ADD API KEY HERE
app_data = {
    'cv_versions': [],
    'current_cv_index': -1
}
saved_jobs = []

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

def parse_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def parse_pdf(file_path):
    text = []
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            extracted = page.extract_text()
            if extracted:
                text.append(extracted)
    return "\n".join(text)

def parse_document(file_path):
    if file_path.endswith('.pptx'):
        return parse_pptx(file_path)
    elif file_path.endswith('.pdf'):
        return parse_pdf(file_path)
    else:
        raise ValueError("Unsupported file type")

def extract_cv_info(cv_text):
    prompt = f"""Analyze this CV and extract the key information.

CV Content:
{cv_text}

Extract and return as JSON:
{{
    "name": "Full name",
    "email": "Email address",
    "phone": "Phone number",
    "current_role": "Current or target role",
    "skills": ["list of technical and soft skills"],
    "experience": ["brief description of work experience"],
    "education": ["education background"],
    "projects": ["notable projects"],
    "achievements": ["key accomplishments"]
}}

Only include information clearly present. Return Not found for missing fields."""

    message = client.messages.create(
        model="claude-opus-4-20250514",
        max_tokens=1500,
        messages=[{"role": "user", "content": prompt}]
    )
    return message.content[0].text

def extract_profile(slide_text):
    prompt = f"""Analyze this content extracted from presentation/lecture slides/project documents and identify professional information for a CV.

Content:
{slide_text}

Extract and return as JSON with COMPREHENSIVE skill extraction:
{{
    "skills": ["list ALL technical and soft skills - be thorough and specific"],
    "projects": ["list of projects mentioned with brief descriptions"],
    "achievements": ["list of accomplishments"],
    "tools": ["technologies, tools, languages, frameworks, platforms mentioned"],
    "expertise_areas": ["broader domains of knowledge"],
    "technical_skills": ["programming languages, databases, cloud platforms, DevOps tools"],
    "soft_skills": ["communication, leadership, teamwork, problem-solving abilities"],
    "methodologies": ["Agile, Scrum, CI/CD, testing methodologies"]
}}

Be VERY thorough with skill extraction:
- Extract ALL programming languages mentioned
- List ALL tools and technologies
- Identify frameworks and libraries
- Note methodologies and best practices
- Include both explicit and implied skills
- Add soft skills demonstrated through projects/presentations

Return comprehensive lists without being conservative."""

    message = client.messages.create(
        model="claude-opus-4-20250514",
        max_tokens=2000,
        messages=[{"role": "user", "content": prompt}]
    )
    return message.content[0].text

def generate_cv(profile_data, user_info, source_cv=None):
    base_info = ""
    if source_cv:
        base_info = f"\n\nBase CV Information (use this as foundation):\n{source_cv}"
    
    prompt = f"""Create a professional, well-formatted CV based on this extracted data:

Profile Data: {profile_data}

User Information:
- Name: {user_info.get('name', 'Not provided')}
- Email: {user_info.get('email', 'Not provided')}
- Phone: {user_info.get('phone', 'Not provided')}
- Target Role: {user_info.get('target_role', 'Professional role')}
{base_info}

Create a complete CV with these sections:
1. HEADER - Name and contact info
2. PROFESSIONAL SUMMARY - 2-3 compelling sentences
3. SKILLS - Organized by category
4. PROJECTS & EXPERIENCE - With impact statements
5. ACHIEVEMENTS - Quantified where possible

Use strong action verbs. Be specific and impactful.
Format it cleanly with clear section headers."""

    message = client.messages.create(
        model="claude-opus-4-20250514",
        max_tokens=2500,
        messages=[{"role": "user", "content": prompt}]
    )
    return message.content[0].text

def tailor_cv_to_job(generic_cv, job_info):
    prompt = f"""Tailor this CV to match the job description. Keep the same format but emphasize relevant skills and experience.

Generic CV:
{generic_cv}

Job Details:
- Job Title: {job_info.get('title', 'Not specified')}
- Company: {job_info.get('company', 'Not specified')}
- Job Description: {job_info.get('description', 'Not provided')}

Instructions:
1. Keep the same CV structure and format
2. Reorder and emphasize skills that match the job requirements
3. Highlight relevant experience and projects
4. Adjust the professional summary to align with the role
5. Maintain all factual information from the original CV
6. Use keywords from the job description naturally

Return the complete tailored CV in the same format."""

    message = client.messages.create(
        model="claude-opus-4-20250514",
        max_tokens=3000,
        messages=[{"role": "user", "content": prompt}]
    )
    return message.content[0].text

def generate_cover_letter(cv_data, job_info):
    prompt = f"""Write a tailored, compelling cover letter.

Candidate Profile:
{cv_data}

Job Details:
- Job Title: {job_info.get('title', 'Not specified')}
- Company: {job_info.get('company', 'Not specified')}
- Job Description: {job_info.get('description', 'Not provided')}

Write a cover letter that:
1. Opens with an engaging hook (not I am writing to apply...)
2. Connects specific skills/projects to job requirements
3. Shows enthusiasm for the company/role
4. Includes a concrete achievement example
5. Ends with a confident call to action

Length: 3-4 paragraphs. Tone: Professional but personable."""

    message = client.messages.create(
        model="claude-opus-4-20250514",
        max_tokens=1500,
        messages=[{"role": "user", "content": prompt}]
    )
    return message.content[0].text

def analyze_job_match(cv_data, job_info):
    prompt = f"""Analyze how well this candidate matches the job and provide a detailed compatibility assessment.

Candidate CV:
{cv_data}

Job Details:
- Job Title: {job_info.get('title', 'Not specified')}
- Company: {job_info.get('company', 'Not specified')}
- Job Description: {job_info.get('description', 'Not provided')}

Return ONLY valid JSON (no other text) with this exact structure:
{{
    "match_score": 85,
    "matched_skills": ["skill1", "skill2", "skill3"],
    "missing_skills": ["skill1", "skill2"],
    "matched_experience": ["relevant experience or project"],
    "key_requirements_met": ["requirement 1", "requirement 2"],
    "gaps": ["gap 1", "gap 2"],
    "recommendation": "Brief 1-2 sentence recommendation",
    "compatibility_breakdown": {{
        "technical_match": 90,
        "experience_match": 80,
        "education_match": 85,
        "soft_skills_match": 75
    }}
}}

The match_score should be a number between 60-98 based on how well the candidate fits.
List 3-6 matched skills and 2-4 missing skills.
Include detailed compatibility breakdown."""

    message = client.messages.create(
        model="claude-opus-4-20250514",
        max_tokens=800,
        messages=[{"role": "user", "content": prompt}]
    )
    
    response_text = message.content[0].text.strip()
    
    try:
        start = response_text.find('{')
        end = response_text.rfind('}') + 1
        if start != -1 and end > start:
            json_str = response_text[start:end]
            return json.loads(json_str)
    except:
        pass
    
    return {
        "match_score": 75,
        "matched_skills": ["Communication", "Problem Solving", "Technical Skills"],
        "missing_skills": ["Specific tools may vary"],
        "matched_experience": ["Relevant projects"],
        "key_requirements_met": ["Basic qualifications"],
        "gaps": ["Some specific tools"],
        "recommendation": "Good potential match. Highlight relevant experience in your application.",
        "compatibility_breakdown": {
            "technical_match": 75,
            "experience_match": 70,
            "education_match": 80,
            "soft_skills_match": 75
        }
    }

def search_recommended_jobs(cv_data, user_skills):
    prompt = f"""Based on this candidate profile, recommend relevant job positions they should consider.

Candidate CV Summary:
{cv_data[:1500]}

Candidate Skills:
{', '.join(user_skills[:50])}

Generate 8-10 realistic job recommendations. Return ONLY valid JSON:
{{
    "recommended_jobs": [
        {{
            "title": "Job Title",
            "company": "Company Name",
            "description": "Brief 2-3 sentence job description focusing on key responsibilities and requirements",
            "match_score": 85,
            "key_skills": ["skill1", "skill2", "skill3"],
            "reason": "Why this is a good match (1 sentence)"
        }}
    ]
}}

Make recommendations realistic and diverse across different companies and seniority levels."""

    message = client.messages.create(
        model="claude-opus-4-20250514",
        max_tokens=2000,
        messages=[{"role": "user", "content": prompt}]
    )
    
    response_text = message.content[0].text.strip()
    
    try:
        start = response_text.find('{')
        end = response_text.rfind('}') + 1
        if start != -1 and end > start:
            json_str = response_text[start:end]
            return json.loads(json_str)
    except:
        pass
    
    return {"recommended_jobs": []}

def save_cv_version(cv_content, description="Manual update"):
    version = {
        'cv': cv_content,
        'description': description,
        'timestamp': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
        'version_number': len(app_data.get('cv_versions', [])) + 1
    }
    
    if 'cv_versions' not in app_data:
        app_data['cv_versions'] = []
    
    app_data['cv_versions'].append(version)
    app_data['current_cv_index'] = len(app_data['cv_versions']) - 1
    app_data['cv'] = cv_content

@app.route('/')
def index():
    return render_template('index_new.html')

@app.route('/upload', methods=['POST'])
def upload():
    has_cv = 'cv_file' in request.files and request.files['cv_file'].filename != ''
    has_slides = 'slides_file' in request.files and request.files['slides_file'].filename != ''
    
    if not has_cv and not has_slides:
        flash('Please upload at least your CV or slides', 'error')
        return redirect(url_for('index'))
    
    user_info = {
        'name': request.form.get('name'),
        'email': request.form.get('email'),
        'phone': request.form.get('phone'),
        'target_role': request.form.get('target_role')
    }
    
    cv_text = None
    profile_data = None
    
    try:
        if has_cv:
            cv_file = request.files['cv_file']
            if allowed_file(cv_file.filename):
                filename = secure_filename(cv_file.filename)
                file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                cv_file.save(file_path)
                
                cv_text = parse_document(file_path)
                cv_info = extract_cv_info(cv_text)
                
                try:
                    cv_data = json.loads(cv_info)
                    if not user_info['name'] or user_info['name'] == '':
                        user_info['name'] = cv_data.get('name', '')
                    if not user_info['email'] or user_info['email'] == '':
                        user_info['email'] = cv_data.get('email', '')
                    if not user_info['phone'] or user_info['phone'] == '':
                        user_info['phone'] = cv_data.get('phone', '')
                    if not user_info['target_role'] or user_info['target_role'] == '':
                        user_info['target_role'] = cv_data.get('current_role', '')
                except:
                    pass
                
                os.remove(file_path)
        
        if has_slides:
            slides_file = request.files['slides_file']
            if allowed_file(slides_file.filename):
                filename = secure_filename(slides_file.filename)
                file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                slides_file.save(file_path)
                
                slide_text = parse_document(file_path)
                profile_data = extract_profile(slide_text)
                
                os.remove(file_path)
        
        cv = generate_cv(profile_data or "{}", user_info, cv_text)
        
        app_data['cv_versions'] = []
        save_cv_version(cv, "Initial CV generation")
        
        app_data['profile'] = profile_data
        app_data['user_info'] = user_info
        app_data['source_cv'] = cv_text
        app_data['enhancements'] = []
        
        return redirect(url_for('dashboard'))
        
    except Exception as e:
        flash(f'Error processing file: {str(e)}', 'error')
        return redirect(url_for('index'))

@app.route('/dashboard')
def dashboard():
    if 'cv' not in app_data:
        return redirect(url_for('index'))
    
    all_skills = {
        'technical': set(),
        'soft': set(),
        'tools': set(),
        'methodologies': set(),
        'all': set()
    }
    
    # Get skills from original profile
    if 'profile' in app_data and app_data['profile']:
        try:
            profile = json.loads(app_data['profile'])
            # Add all skill types to 'all' category
            all_skills['all'].update(profile.get('skills', []))
            all_skills['technical'].update(profile.get('technical_skills', []))
            all_skills['soft'].update(profile.get('soft_skills', []))
            all_skills['tools'].update(profile.get('tools', []))
            all_skills['methodologies'].update(profile.get('methodologies', []))
            all_skills['all'].update(profile.get('technical_skills', []))
            all_skills['all'].update(profile.get('soft_skills', []))
            all_skills['all'].update(profile.get('tools', []))
            all_skills['all'].update(profile.get('methodologies', []))
        except Exception as e:
            print(f"Error parsing profile: {e}")
            print(f"Profile data: {app_data.get('profile', 'None')}")
    
    # Get skills from enhancements
    if 'enhancements' in app_data:
        for enhancement in app_data['enhancements']:
            try:
                profile = json.loads(enhancement['profile'])
                all_skills['all'].update(profile.get('skills', []))
                all_skills['technical'].update(profile.get('technical_skills', []))
                all_skills['soft'].update(profile.get('soft_skills', []))
                all_skills['tools'].update(profile.get('tools', []))
                all_skills['methodologies'].update(profile.get('methodologies', []))
                all_skills['all'].update(profile.get('technical_skills', []))
                all_skills['all'].update(profile.get('soft_skills', []))
                all_skills['all'].update(profile.get('tools', []))
                all_skills['all'].update(profile.get('methodologies', []))
            except Exception as e:
                print(f"Error parsing enhancement: {e}")
    
    # Combine all unique skills
    combined_skills = all_skills['all']
    
    # Debug output
    print(f"Total skills found: {len(combined_skills)}")
    print(f"Technical: {len(all_skills['technical'])}")
    print(f"Soft: {len(all_skills['soft'])}")
    print(f"Tools: {len(all_skills['tools'])}")
    print(f"Methodologies: {len(all_skills['methodologies'])}")
    
    return render_template('dashboard.html',
                         user_info=app_data.get('user_info'),
                         skills=sorted(list(combined_skills)),
                         skills_by_category={
                             'technical': sorted(list(all_skills['technical'])),
                             'soft': sorted(list(all_skills['soft'])),
                             'tools': sorted(list(all_skills['tools'])),
                             'methodologies': sorted(list(all_skills['methodologies']))
                         },
                         enhancement_count=len(app_data.get('enhancements', [])),
                         job_count=len(saved_jobs),
                         cv_version=app_data.get('current_cv_index', 0) + 1,
                         total_versions=len(app_data.get('cv_versions', [])))

@app.route('/enhance-profile')
def enhance_profile():
    if 'cv' not in app_data:
        flash('Please generate your base CV first', 'error')
        return redirect(url_for('index'))
    return render_template('enhance_profile.html', 
                         user_info=app_data.get('user_info'),
                         enhancements=app_data.get('enhancements', []))

@app.route('/enhance-profile/upload', methods=['POST'])
def upload_enhancement():
    if 'cv' not in app_data:
        flash('Please generate your base CV first', 'error')
        return redirect(url_for('index'))
    
    if 'file' not in request.files:
        flash('No file selected', 'error')
        return redirect(url_for('enhance_profile'))
    
    file = request.files['file']
    if file.filename == '':
        flash('No file selected', 'error')
        return redirect(url_for('enhance_profile'))
    
    if file and allowed_file(file.filename):
        try:
            filename = secure_filename(file.filename)
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(file_path)
            
            content = parse_document(file_path)
            profile_data = extract_profile(content)
            
            if 'enhancements' not in app_data:
                app_data['enhancements'] = []
            
            app_data['enhancements'].append({
                'filename': file.filename,
                'content': content,
                'profile': profile_data,
                'date_added': datetime.now().strftime('%Y-%m-%d %H:%M')
            })
            
            regenerate_enhanced_cv()
            
            os.remove(file_path)
            
            flash(f'Successfully added {file.filename} to your profile!', 'success')
            return redirect(url_for('enhance_profile'))
            
        except Exception as e:
            flash(f'Error processing file: {str(e)}', 'error')
            return redirect(url_for('enhance_profile'))
    
    flash('Invalid file type. Please upload PDF or PPTX.', 'error')
    return redirect(url_for('enhance_profile'))

@app.route('/enhance-profile/remove/<int:index>', methods=['POST'])
def remove_enhancement(index):
    if 'enhancements' in app_data and 0 <= index < len(app_data['enhancements']):
        removed = app_data['enhancements'].pop(index)
        regenerate_enhanced_cv()
        return jsonify({'success': True, 'message': f'Removed {removed["filename"]}'})
    
    return jsonify({'success': False, 'message': 'Enhancement not found'}), 404

def regenerate_enhanced_cv():
    all_profiles = []
    
    if 'profile' in app_data and app_data['profile']:
        all_profiles.append(app_data['profile'])
    
    if 'enhancements' in app_data:
        for enhancement in app_data['enhancements']:
            all_profiles.append(enhancement['profile'])
    
    combined_profile = {
        'skills': [],
        'projects': [],
        'achievements': [],
        'tools': [],
        'expertise_areas': [],
        'technical_skills': [],
        'soft_skills': [],
        'methodologies': []
    }
    
    for profile_str in all_profiles:
        try:
            profile = json.loads(profile_str)
            for key in combined_profile.keys():
                if key in profile:
                    combined_profile[key].extend(profile[key])
        except:
            pass
    
    for key in combined_profile.keys():
        seen = set()
        unique_list = []
        for item in combined_profile[key]:
            if item not in seen:
                seen.add(item)
                unique_list.append(item)
        combined_profile[key] = unique_list
    
    enhanced_cv = generate_cv(
        json.dumps(combined_profile),
        app_data.get('user_info', {}),
        app_data.get('source_cv')
    )
    
    enhancement_files = [e['filename'] for e in app_data.get('enhancements', [])]
    description = f"Enhanced with: {', '.join(enhancement_files[-3:])}" if enhancement_files else "Profile enhancement"
    save_cv_version(enhanced_cv, description)
    
    app_data['combined_profile'] = combined_profile

@app.route('/results')
def results():
    if 'cv' not in app_data:
        return redirect(url_for('index'))
    return render_template('results.html', 
                         cv=app_data.get('cv'),
                         profile=app_data.get('profile'),
                         user_info=app_data.get('user_info'))

@app.route('/jobs')
def jobs_manager():
    return render_template('jobs.html', jobs=saved_jobs)

@app.route('/jobs/add', methods=['POST'])
def add_job():
    job = {
        'id': len(saved_jobs) + 1,
        'title': request.form.get('job_title'),
        'company': request.form.get('company'),
        'description': request.form.get('job_description'),
        'date_added': datetime.now().strftime('%Y-%m-%d %H:%M'),
        'applied': False
    }
    saved_jobs.append(job)
    flash('Job description saved successfully!', 'success')
    return redirect(url_for('jobs_manager'))

@app.route('/jobs/delete/<int:job_id>', methods=['POST'])
def delete_job(job_id):
    global saved_jobs
    saved_jobs = [j for j in saved_jobs if j['id'] != job_id]
    return jsonify({'success': True})

@app.route('/jobs/recommendations')
def job_recommendations():
    if 'cv' not in app_data:
        flash('Please generate your base CV first', 'error')
        return redirect(url_for('index'))
    
    all_skills = set()
    if 'profile' in app_data and app_data['profile']:
        try:
            profile = json.loads(app_data['profile'])
            all_skills.update(profile.get('skills', []))
            all_skills.update(profile.get('technical_skills', []))
            all_skills.update(profile.get('tools', []))
        except:
            pass
    
    if 'enhancements' in app_data:
        for enhancement in app_data['enhancements']:
            try:
                profile = json.loads(enhancement['profile'])
                all_skills.update(profile.get('skills', []))
                all_skills.update(profile.get('technical_skills', []))
                all_skills.update(profile.get('tools', []))
            except:
                pass
    
    recommendations_data = search_recommended_jobs(app_data['cv'], list(all_skills))
    recommended_jobs = recommendations_data.get('recommended_jobs', [])
    
    recommended_jobs.sort(key=lambda x: x.get('match_score', 0), reverse=True)
    
    return render_template('job_recommendations.html', 
                         recommended_jobs=recommended_jobs,
                         user_skills=sorted(list(all_skills)))

@app.route('/jobs/add-recommended', methods=['POST'])
def add_recommended_job():
    data = request.get_json()
    
    job = {
        'id': len(saved_jobs) + 1,
        'title': data.get('title'),
        'company': data.get('company'),
        'description': data.get('description'),
        'date_added': datetime.now().strftime('%Y-%m-%d %H:%M'),
        'applied': False,
        'match_score': data.get('match_score', 0)
    }
    saved_jobs.append(job)
    
    return jsonify({'success': True, 'message': 'Job added to your list!'})

@app.route('/skills/delete', methods=['POST'])
def delete_skill():
    data = request.get_json()
    skill_to_remove = data.get('skill')
    
    if not skill_to_remove:
        return jsonify({'success': False, 'message': 'No skill specified'}), 400
    
    if 'profile' in app_data and app_data['profile']:
        try:
            profile = json.loads(app_data['profile'])
            for key in ['skills', 'technical_skills', 'soft_skills', 'tools', 'methodologies']:
                if key in profile and skill_to_remove in profile[key]:
                    profile[key].remove(skill_to_remove)
            app_data['profile'] = json.dumps(profile)
        except:
            pass
    
    if 'enhancements' in app_data:
        for enhancement in app_data['enhancements']:
            try:
                profile = json.loads(enhancement['profile'])
                for key in ['skills', 'technical_skills', 'soft_skills', 'tools', 'methodologies']:
                    if key in profile and skill_to_remove in profile[key]:
                        profile[key].remove(skill_to_remove)
                enhancement['profile'] = json.dumps(profile)
            except:
                pass
    
    regenerate_enhanced_cv()
    
    return jsonify({'success': True, 'message': f'Removed {skill_to_remove}'})

@app.route('/jobs/apply/<int:job_id>')
def apply_to_job(job_id):
    if 'cv' not in app_data:
        flash('Please generate your base CV first', 'error')
        return redirect(url_for('index'))
    
    job = next((j for j in saved_jobs if j['id'] == job_id), None)
    if not job:
        flash('Job not found', 'error')
        return redirect(url_for('jobs_manager'))
    
    tailored_cv = tailor_cv_to_job(app_data['cv'], job)
    letter = generate_cover_letter(tailored_cv, job)
    match_result = analyze_job_match(tailored_cv, job)
    
    job['applied'] = True
    
    save_cv_version(tailored_cv, f"Tailored for {job['company']} - {job['title']}")
    
    return render_template('application.html',
                         cv=tailored_cv,
                         cover_letter=letter,
                         job_info=job,
                         match_result=match_result)

@app.route('/cv-versions')
def cv_versions():
    if 'cv' not in app_data:
        return redirect(url_for('index'))
    
    versions = app_data.get('cv_versions', [])
    return render_template('cv_versions.html', 
                         versions=versions,
                         current_index=app_data.get('current_cv_index', -1))

@app.route('/cv-versions/switch/<int:index>')
def switch_cv_version(index):
    versions = app_data.get('cv_versions', [])
    if 0 <= index < len(versions):
        app_data['current_cv_index'] = index
        app_data['cv'] = versions[index]['cv']
        flash(f'Switched to CV version {index + 1}', 'success')
    
    return redirect(url_for('cv_versions'))

@app.route('/cover-letter', methods=['GET', 'POST'])
def cover_letter():
    if 'cv' not in app_data:
        return redirect(url_for('index'))
    
    if request.method == 'POST':
        job_info = {
            'title': request.form.get('job_title'),
            'company': request.form.get('company'),
            'description': request.form.get('job_description')
        }
        
        letter = generate_cover_letter(app_data['cv'], job_info)
        match_result = analyze_job_match(app_data['cv'], job_info)
        
        app_data['cover_letter'] = letter
        app_data['job_info'] = job_info
        app_data['match_result'] = match_result
        
        return render_template('cover_letter.html',
                             cover_letter=letter,
                             job_info=job_info,
                             match_result=match_result)
    
    if 'cover_letter' in app_data:
        return render_template('cover_letter.html',
                             cover_letter=app_data['cover_letter'],
                             job_info=app_data['job_info'],
                             match_result=app_data.get('match_result', {}))
    
    return render_template('cover_letter_form.html')

@app.route('/new-cover-letter')
def new_cover_letter():
    if 'cover_letter' in app_data:
        del app_data['cover_letter']
    if 'job_info' in app_data:
        del app_data['job_info']
    if 'match_result' in app_data:
        del app_data['match_result']
    return redirect(url_for('cover_letter'))

@app.route('/reset')
def reset():
    app_data.clear()
    global saved_jobs
    saved_jobs = []
    return redirect(url_for('index'))

@app.route('/debug')
def debug():
    """Debug route to see what data is stored"""
    debug_info = {
        'has_cv': 'cv' in app_data,
        'has_profile': 'profile' in app_data,
        'profile_data': app_data.get('profile', 'None'),
        'enhancements_count': len(app_data.get('enhancements', [])),
        'cv_versions': len(app_data.get('cv_versions', [])),
        'user_info': app_data.get('user_info', {})
    }
    
    # Try to parse profile
    if app_data.get('profile'):
        try:
            profile = json.loads(app_data['profile'])
            debug_info['parsed_profile'] = {
                'skills_count': len(profile.get('skills', [])),
                'technical_skills_count': len(profile.get('technical_skills', [])),
                'tools_count': len(profile.get('tools', [])),
                'soft_skills_count': len(profile.get('soft_skills', [])),
                'methodologies_count': len(profile.get('methodologies', [])),
                'sample_skills': profile.get('skills', [])[:5],
                'sample_technical': profile.get('technical_skills', [])[:5]
            }
        except Exception as e:
            debug_info['parse_error'] = str(e)
    
    return jsonify(debug_info)

if __name__ == '__main__':
    app.run(debug=True)